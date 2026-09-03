#!/usr/bin/env python3
"""
survey_flowchart.py

Generates a CONSORT-style participant flow diagram as a PowerPoint slide.
Edit the COUNTS dict below to update numbers for each wave of data collection.

Output: survey_flowchart.pptx  (same directory as this script)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Editable counts ────────────────────────────────────────────────────────────
COUNTS = {
    "total_navigated":          471,   # rows in CSV (all who navigated to survey)
    "screener_incomplete":       25,   # did not finish screener (Screener_complete ≠ Complete)
    "screener_complete":        446,   # = total_navigated - screener_incomplete
    "blank_eligibility":          2,   # screener-complete but blank eligibility answer
    "ineligible":               291,   # do not provide prenatal/delivery/postpartum care
    "eligible":                 153,   # advanced to full survey
    "partial_completers":        31,   # Survey_complete = Incomplete
    "full_completers":          122,   # Survey_complete = Complete (final analytic sample)
    # Ineligible demographics completion (from Inel_demo_complete column)
    "inel_demo_complete":        None,  # fill in once known (e.g. 246)
    "inel_demo_incomplete":      None,  # fill in once known (e.g. 21)
}

OUT_PATH = Path(__file__).parent / "survey_flowchart.pptx"

# ── Layout constants ───────────────────────────────────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

BLUE   = RGBColor(0x1F, 0x49, 0x7D)   # dark blue — main boxes
GREY   = RGBColor(0xD9, 0xD9, 0xD9)   # light grey — exclusion / side boxes
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
LGREEN = RGBColor(0xE2, 0xEF, 0xDA)   # light green — final sample box

FONT   = "Calibri"

def emu(inches):
    return Inches(inches)


def add_box(slide, left, top, width, height,
            text, font_size=11, bold=False,
            fill=BLUE, text_color=WHITE, align=PP_ALIGN.CENTER,
            border_color=None):
    shape = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    fill_obj = shape.fill
    fill_obj.solid()
    fill_obj.fore_color.rgb = fill
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    if border_color:
        ln = shape.line
        ln.color.rgb = border_color
        ln.width = Pt(0.75)
    return shape


def add_arrow(slide, x, y1, y2, color=BLACK):
    """Vertical arrow from (x, y1) down to (x, y2)."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        emu(x), emu(y1), emu(x), emu(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


def add_horiz_arrow(slide, x1, x2, y, color=BLACK):
    """Horizontal arrow from (x1, y) to (x2, y)."""
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,
        emu(x1), emu(y), emu(x2), emu(y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


def build_flowchart(counts):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Title ──────────────────────────────────────────────────────────────────
    add_box(slide, 0.3, 0.1, 9.4, 0.45,
            "Figure 1. Survey Participant Flow",
            font_size=13, bold=True, fill=WHITE, text_color=BLACK)

    # ── Main column x-centre = 3.5, width = 3.2 ───────────────────────────────
    cx, cw = 3.4, 3.2

    # Row positions (top of each box)
    row = [0.65, 1.55, 2.45, 3.50, 4.65, 5.55]
    bh  = 0.65   # box height

    # Box 1: Total navigated
    add_box(slide, cx, row[0], cw, bh,
            f"{counts['total_navigated']:,} navigated to survey\n(via Z-link or QR code)",
            fill=BLUE, text_color=WHITE, bold=False)

    # Arrow 1→2
    add_arrow(slide, cx + cw/2, row[0]+bh, row[1])

    # Box 2: Screener complete
    add_box(slide, cx, row[1], cw, bh,
            f"{counts['screener_complete']:,} completed screening questions",
            fill=BLUE, text_color=WHITE)

    # Side box: screener exclusion
    exc1_txt = f"{counts['screener_incomplete']:,} did not complete\nscreening questions"
    add_box(slide, 7.0, row[1]+0.05, 2.6, 0.55,
            exc1_txt, font_size=10, fill=GREY, text_color=BLACK,
            align=PP_ALIGN.LEFT, border_color=BLACK)
    add_horiz_arrow(slide, cx+cw, 7.0, row[1]+bh/2)

    # Arrow 2→3
    add_arrow(slide, cx + cw/2, row[1]+bh, row[2])

    # Box 3: Eligibility determined
    add_box(slide, cx, row[2], cw, 0.45,
            "Eligibility determined from screening",
            fill=BLUE, text_color=WHITE, font_size=10)

    add_arrow(slide, cx + cw/2, row[2]+0.45, row[3])

    # Row 3: split — ineligible (left side box) and eligible (right centre)
    # Ineligible box
    inel_lines = [f"{counts['ineligible']:,} Ineligible",
                  "Directed to demographics\nsection only"]
    if counts["inel_demo_complete"]:
        inel_lines.append(f"  {counts['inel_demo_complete']:,} completed & submitted")
    if counts["inel_demo_incomplete"]:
        inel_lines.append(f"  {counts['inel_demo_incomplete']:,} did not submit")
    add_box(slide, 0.3, row[3], 2.8, bh+0.15,
            "\n".join(inel_lines),
            font_size=10, fill=GREY, text_color=BLACK,
            align=PP_ALIGN.LEFT, border_color=BLACK)

    # Eligible box
    add_box(slide, cx, row[3], cw, bh,
            f"{counts['eligible']:,} Eligible\nAdvanced to full survey",
            fill=BLUE, text_color=WHITE, bold=False)

    # Horizontal connector from centre to ineligible
    add_horiz_arrow(slide, cx, 0.3+2.8, row[3]+bh/2)

    # Exclusion 2 — blank eligibility (small note on right)
    if counts["blank_eligibility"]:
        add_box(slide, 7.0, row[3]+0.05, 2.6, 0.55,
                f"{counts['blank_eligibility']:,} excluded: blank\neligibility response",
                font_size=10, fill=GREY, text_color=BLACK,
                align=PP_ALIGN.LEFT, border_color=BLACK)
        add_horiz_arrow(slide, cx+cw, 7.0, row[3]+bh/2)

    # Arrow eligible → completers
    add_arrow(slide, cx + cw/2, row[3]+bh, row[4])

    # Partial completer note (right side)
    add_box(slide, 7.0, row[4]+0.05, 2.6, 0.55,
            f"{counts['partial_completers']:,} initiated but did\nnot submit",
            font_size=10, fill=GREY, text_color=BLACK,
            align=PP_ALIGN.LEFT, border_color=BLACK)
    add_horiz_arrow(slide, cx+cw, 7.0, row[4]+bh/2)

    # Box 5: Final analytic sample
    add_box(slide, cx, row[4], cw, bh,
            f"{counts['full_completers']:,} completed & submitted\n(final analytic sample)",
            fill=LGREEN, text_color=BLACK, bold=True,
            border_color=BLUE)

    # ── Footnote ───────────────────────────────────────────────────────────────
    note = (
        'Note: "Completed" = navigated through the entire survey and submitted '
        "(does not require answering every question). "
        f"Of the {counts['partial_completers']} who did not submit: "
        "see resp_complete_detail sheet for question-level completion rates."
    )
    add_box(slide, 0.3, 6.7, 9.4, 0.72,
            note, font_size=9, fill=WHITE, text_color=BLACK,
            align=PP_ALIGN.LEFT)

    prs.save(OUT_PATH)
    print(f"Flowchart saved: {OUT_PATH}")


if __name__ == "__main__":
    build_flowchart(COUNTS)
